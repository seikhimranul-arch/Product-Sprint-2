from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette (FintLer brand) ──
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_SURFACE = RGBColor(0x0A, 0x0A, 0x0F)
CARD_BG = RGBColor(0x14, 0x14, 0x1E)
TEAL = RGBColor(0x00, 0xDF, 0xC1)
TEAL_DIM = RGBColor(0x00, 0xAA, 0x99)
PURPLE = RGBColor(0xBB, 0x86, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xAD)
ERROR = RGBColor(0xFF, 0x6B, 0x6B)
GOLD = RGBColor(0xFF, 0xD7, 0x00)

def add_bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_orb(slide, left, top, size, color, opacity=0.15):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.fore_color.brightness = opacity
    line = shape.line
    line.fill.background()
    shape.shadow.inherit = False
    # Move to back
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=GRAY, bullet_color=TEAL):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return txBox

def add_card(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    line = shape.line
    line.color.rgb = RGBColor(0x22, 0x22, 0x30)
    line.width = Pt(1)
    return shape

def add_stat_card(slide, left, top, width, height, label, value, accent_color=TEAL):
    add_card(slide, left, top, width, height)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.2), width - Inches(0.6), Inches(0.4),
                label, 11, GRAY)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.6), width - Inches(0.6), Inches(0.6),
                value, 24, accent_color, bold=True)

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_orb(slide, Inches(-2), Inches(-1), Inches(8), TEAL, 0.12)
add_orb(slide, Inches(8), Inches(3), Inches(6), PURPLE, 0.10)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
            "FintLer", 48, TEAL, bold=True)
add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
            "Cure Salary Day Anxiety. See Exactly Where Your Money Vanishes.", 24, WHITE)
add_textbox(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.5),
            "A Behavioral Financial Clarity Engine for India's UPI Generation", 16, GRAY)
add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
            "Product Deck  •  July 2026", 14, RGBColor(0x60, 0x60, 0x70))

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.7),
            "The Problem: Salary Day Spending Anxiety", 32, WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(6), Inches(0.5),
            "India's young workforce is drowning in transactions, starved of clarity.", 16, GRAY)

add_stat_card(slide, Inches(0.8), Inches(2.0), Inches(2.8), Inches(1.3), "Gen Z financially insecure", "54%", ERROR)
add_stat_card(slide, Inches(3.9), Inches(2.0), Inches(2.8), Inches(1.3), "Millennials financially insecure", "46%", ERROR)
add_stat_card(slide, Inches(7.0), Inches(2.0), Inches(2.8), Inches(1.3), "Live paycheque-to-paycheque", "48%", ERROR)
add_stat_card(slide, Inches(10.1), Inches(2.0), Inches(2.8), Inches(1.3), "Report daily money anxiety", "80%", GOLD)

add_bullet_list(slide, Inches(0.8), Inches(3.8), Inches(5.5), Inches(3.2), [
    "India's $148B fintech market overlooks behavioral clarity",
    "Existing apps: manual entry (abandoned), judgmental tone, or bank-locked",
    "UPI processes 19.6B transactions/month — but no one explains the story",
    "57.6% budget — but only half use tracking apps effectively",
    "Western budgeting models don't fit India's UPI-first, multi-bank reality",
], 15)

add_bullet_list(slide, Inches(6.8), Inches(3.8), Inches(5.5), Inches(3.2), [
    "Avg Indian salaried millennial: ₹42,500/mo UPI spend",
    "4+ streaming subscriptions, 3+ UPI apps, 2+ bank accounts",
    "EMI deductions cause mid-month balance anxiety",
    "Salary Day = relief + overspend guilt cycle repeated monthly",
    "No tool connects the dots between UPI, subscriptions, EMI & lifestyle",
], 15, PURPLE)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_orb(slide, Inches(9), Inches(4), Inches(5), TEAL, 0.10)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "FintLer: Your Financial Clarity Engine", 32, WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
            "Zero manual entry. Zero judgment. Pure behavioral clarity.", 16, GRAY)

# Three core pillars
for i, (title, desc, icon, accent) in enumerate([
    ("Zero-Data Entry", "Gmail integration parses bank alerts automatically. Works across HDFC, ICICI, SBI, Axis, Kotak, Yes Bank. No manual logging = no abandonment.", "🔌", TEAL),
    ("AI Behavioral Engine", "Gemini 2.0 Flash categorizes transactions, detects patterns & generates personality profiles. 95%+ extraction accuracy on Indian bank emails.", "🧠", PURPLE),
    ("Empowering Insights", "Personality labels ('Weekend Warrior'), behavioral triggers ('Late-night food = ₹4,200/mo'), micro-goals ('Try ₹500 Swiggy cap'). No guilt, just clarity.", "🎯", GOLD),
]):
    left = Inches(0.8 + i * 4.1)
    add_card(slide, left, Inches(2.0), Inches(3.8), Inches(4.5))
    add_textbox(slide, left + Inches(0.3), Inches(2.2), Inches(0.6), Inches(0.6), icon, 28, accent)
    add_textbox(slide, left + Inches(0.3), Inches(2.9), Inches(3.3), Inches(0.4), title, 20, accent, bold=True)
    add_textbox(slide, left + Inches(0.3), Inches(3.4), Inches(3.3), Inches(3.0), desc, 14, GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — Market Opportunity
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "India Fintech — The Opportunity", 32, WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
            "The fastest-growing fintech market globally, with a behavioral clarity gap at its center.", 16, GRAY)

add_stat_card(slide, Inches(0.8), Inches(2.0), Inches(3.0), Inches(1.3), "India Fintech Market (2026)", "$148B", TEAL)
add_stat_card(slide, Inches(4.1), Inches(2.0), Inches(3.0), Inches(1.3), "Projected by 2033", "$867B", TEAL)
add_stat_card(slide, Inches(7.4), Inches(2.0), Inches(3.0), Inches(1.3), "CAGR (2026–2033)", "28.7%", GOLD)
add_stat_card(slide, Inches(10.7), Inches(2.0), Inches(2.0), Inches(1.3), "Personal Finance Mkt", "$44.5M", PURPLE)

add_bullet_list(slide, Inches(0.8), Inches(3.8), Inches(5.6), Inches(3.2), [
    "UPI: 491M users, 65M merchants, 19.6B monthly transactions",
    "Digital payments = 42.87% of fintech market share",
    "Mobile apps = 67.83% of user interfaces (fastest growing at 18.39%)",
    "Retail users = 66.24% of market — individual consumers dominate",
    "Neobanking fastest-growing segment at 19.64% CAGR",
], 14)

add_bullet_list(slide, Inches(6.8), Inches(3.8), Inches(5.5), Inches(3.2), [
    "Gen Z/millennials: 600M+ target demographic in India",
    "48% Gen Z financially insecure — actively seeking solutions",
    "80% report money anxiety — massive unfilled need",
    "57.6% already budget — but need behavioral 'why', not just 'what'",
    "Account Aggregator framework unlocking data access",
], 14, PURPLE)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — Competitive Landscape
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Competitive Landscape — The White Space", 32, WHITE, bold=True)

segments = [
    ("Payments", "PhonePe, GPay, Paytm", "Transaction rails only. No financial clarity, no behavioral insights.", TEAL),
    ("Budgeting", "CRED, ET Money", "Reward/UPS-driven. Generic categories. No bank-agnostic tracking.", PURPLE),
    ("Neobanking", "Fi, Jupiter, Niyo", "Partner-bank-locked. Only work with 1-2 banks. Can't see full picture.", GOLD),
    ("Wealth", "Groww, Zerodha", "Investment-only. No expense management, no behavioral analysis.", ERROR),
    ("Credit", "Moneyview, CASHe", "Lending-focused. Not financial wellness. High APR products.", RGBColor(0xFF, 0x99, 0x00)),
]

for i, (segment, players, gap, accent) in enumerate(segments):
    y = Inches(1.4 + i * 1.15)
    add_card(slide, Inches(0.8), y, Inches(3.0), Inches(0.95))
    add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(2.5), Inches(0.4), segment, 18, accent, bold=True)
    add_textbox(slide, Inches(1.0), y + Inches(0.5), Inches(2.5), Inches(0.3), players, 11, GRAY)

    add_card(slide, Inches(4.0), y, Inches(4.5), Inches(0.95))
    add_textbox(slide, Inches(4.2), y + Inches(0.15), Inches(4.1), Inches(0.7), gap, 12, GRAY)

# FintLer highlight
add_card(slide, Inches(8.7), Inches(1.4), Inches(3.8), Inches(5.3))
add_textbox(slide, Inches(9.0), Inches(1.6), Inches(3.2), Inches(0.5),
            "★ FintLer Differentiators", 18, TEAL, bold=True)
add_bullet_list(slide, Inches(9.0), Inches(2.2), Inches(3.2), Inches(4.2), [
    "Bank-agnostic: works across ALL banks via Gmail",
    "Zero manual entry: parse, don't type",
    "Behavioral AI: personality, triggers, micro-goals",
    "Zero judgment tone: empower, don't shame",
    "Real-time: Pub/Sub webhooks for instant updates",
    "Consent-first: read-only, privacy-respecting",
    "Indian-first: UPI, ₹, ₹42K patterns, multi-bank",
], 13, WHITE)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — User Persona
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Target Persona: Indian Salary Millennial", 32, WHITE, bold=True)

add_card(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.0))
add_textbox(slide, Inches(1.1), Inches(1.7), Inches(5.0), Inches(0.4),
            "👤  Arjun — The Salary Millennial", 22, TEAL, bold=True)
add_bullet_list(slide, Inches(1.1), Inches(2.3), Inches(5.0), Inches(4.0), [
    "Age: 28, Software Engineer in Bangalore",
    "Income: ₹18 LPA | Rent: ₹28K | EMI: ₹8,500",
    "Has: HDFC + ICICI accounts, 3 UPI apps, 4 subscriptions",
    "Monthly spend: ~₹45K — 42% on food transport",
    "Anxious about: mid-month balance, EMI auto-debit, savings rate",
    "Currently: checks 3 bank apps + CRED manually every week",
    "Wants: 'One place to see where it all goes — without judgment'",
    "Tech: Uses 8+ apps daily, trusts Google, hates manual entry",
], 14, WHITE)

add_card(slide, Inches(6.6), Inches(1.5), Inches(5.5), Inches(5.0))
add_textbox(slide, Inches(6.9), Inches(1.7), Inches(5.0), Inches(0.4),
            "📊  User Behavior Patterns (from demo data)", 22, TEAL, bold=True)
add_bullet_list(slide, Inches(6.9), Inches(2.3), Inches(5.0), Inches(4.0), [
    "Salary credit: ₹85K on 1st of every month",
    "Food delivery: 25+ orders/month, avg ₹650 each",
    "Peak spend: Friday 10 PM–midnight (late-night cravings)",
    "Weekend surge: 62% of discretionary spend Fri–Sun",
    "Subscription bleed: ₹1,216/mo on 4 streaming services",
    "EMI anchor: ₹8,500/mo = 20% of take-home pay",
    "Transfer pattern: ₹2K–3K UPI to friends weekly",
    "Grocery spend: well-controlled at ₹6,500/mo",
], 14, WHITE)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — Product Demo Flow
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "User Journey → Product Experience", 32, WHITE, bold=True)

steps = [
    ("1", "Landing", "3-tier consent + Google OAuth with Gmail readonly scope"),
    ("2", "Sync", "Gmail API scans 90 days of bank alert emails"),
    ("3", "Parse", "Gemini 2.0 extracts amount, merchant, category, date"),
    ("4", "Personality", "AI assigns spending persona + behavioral profile"),
    ("5", "Dashboard", "Bento grid: personality, summary, alert, goal, transactions"),
    ("6", "Real-time", "Pub/Sub webhook triggers instant re-analysis on new email"),
    ("7", "Weekly Goal", "Actionable micro-targets built from behavioral data"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + (i % 4) * 3.2)
    y = Inches(1.5 + (i // 4) * 2.8)
    add_card(slide, x, y, Inches(2.9), Inches(2.3))
    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5),
                num, 28, TEAL, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.5), Inches(0.4),
                title, 18, WHITE, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.15), Inches(2.5), Inches(1.0),
                desc, 12, GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — Architecture
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Technical Architecture", 32, WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
            "Modern, secure, and scalable — built for production.", 16, GRAY)

components = [
    ("Frontend", "React 19 + Vite 8", "Netlify-hosted SPA, Tailwind v4, Framer Motion, code-split vendor chunks", TEAL),
    ("Auth & Data", "Supabase", "Row-Level Security, Google OAuth, Postgres DB, realtime subscriptions", PURPLE),
    ("Email Parsing", "Gemini 2.0 Flash", "Regex-first → AI-fallback pipeline. 95%+ accuracy on Indian bank emails", GOLD),
    ("Behavioral AI", "Gemini 2.0 Flash", "Context-aware personality analysis with ₹-format, UPI pattern recognition", RGBColor(0xFF, 0x99, 0x00)),
    ("Edge Functions", "Deno + Supabase", "6 functions: sync, parse, insights, webhook, debug, statement", ERROR),
    ("Email Sync", "Gmail API", "Read-only scope (gmail.readonly), Pub/Sub webhook for real-time updates", TEAL),
]

for i, (layer, tech, desc, accent) in enumerate(components):
    y = Inches(1.9 + i * 0.9)
    add_card(slide, Inches(0.8), y, Inches(2.5), Inches(0.75))
    add_textbox(slide, Inches(1.0), y + Inches(0.08), Inches(2.1), Inches(0.35), layer, 14, accent, bold=True)
    add_textbox(slide, Inches(1.0), y + Inches(0.38), Inches(2.1), Inches(0.3), tech, 11, GRAY)

    add_card(slide, Inches(3.5), y, Inches(8.8), Inches(0.75))
    add_textbox(slide, Inches(3.7), y + Inches(0.15), Inches(8.4), Inches(0.45), desc, 12, GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 9 — Product Roadmap
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Product Roadmap", 32, WHITE, bold=True)

phases = [
    ("Now (H2 2026)", TEAL, [
        "Private beta: waitlist-gated access",
        "6 major Indian banks supported (HDFC, ICICI, SBI, Axis, Kotak, Yes)",
        "Core behavioral engine: personality + insights + micro-goals",
        "Weekly goal setting with progress tracking",
        "Feedback loop for bank email format training",
    ]),
    ("Next (H1 2027)", PURPLE, [
        "React Native mobile app (67.8% of fintech users mobile-first)",
        "UPI Lite auto-tracking without Gmail",
        "Shared goals for couples/family",
        "Bill negotiation AI agent",
        "Expand to 10+ banks & NBFC email formats",
        "Export to PDF/CSV for tax planning",
    ]),
    ("Future (2028+)", GOLD, [
        "Account Aggregator framework integration",
        "Credit score simulator: 'What-if' on CIBIL",
        "B2B employer wellness program",
        "Financial health API for neobanks & NBFCs",
        "AI coach: conversational financial advisor",
        "Investment recommendation engine (FD → MF → Stocks)",
    ]),
]

for i, (title, accent, items) in enumerate(phases):
    left = Inches(0.8 + i * 4.1)
    add_card(slide, left, Inches(1.4), Inches(3.8), Inches(5.5))
    add_textbox(slide, left + Inches(0.3), Inches(1.6), Inches(3.3), Inches(0.5),
                title, 20, accent, bold=True)
    for j, item in enumerate(items):
        y = Inches(2.3 + j * 0.7)
        add_card(slide, left + Inches(0.3), y, Inches(3.2), Inches(0.55))
        add_textbox(slide, left + Inches(0.5), y + Inches(0.08), Inches(2.8), Inches(0.4),
                    f"→ {item}", 11, WHITE)

# ════════════════════════════════════════════════════════════════
# SLIDE 10 — Metrics & Revenue
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Success Metrics & Revenue Model", 32, WHITE, bold=True)

# Left: Metrics
add_textbox(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
            "Key Metrics (12-Month Targets)", 20, TEAL, bold=True)

metrics = [
    ("Waitlist → Activation", ">40%", "Consent flow isn't blocking"),
    ("Week-1 Retention", ">60%", "Aha moment from first insight"),
    ("Week-4 Retention", ">30%", "Weekly goal habit loop"),
    ("Insight → Action Rate", ">25%", "Goals set after reading AI"),
    ("Gmail Connection Rate", ">70%", "Low-friction OAuth flow"),
    ("NPS Score", ">50", "Product-market fit signal"),
]

for i, (metric, target, reasoning) in enumerate(metrics):
    y = Inches(1.9 + i * 0.75)
    add_card(slide, Inches(0.8), y, Inches(2.2), Inches(0.6))
    add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(1.9), Inches(0.4),
                target, 18, TEAL, bold=True)
    add_card(slide, Inches(3.1), y, Inches(3.6), Inches(0.6))
    add_textbox(slide, Inches(3.3), y + Inches(0.05), Inches(3.3), Inches(0.25),
                metric, 13, WHITE, bold=True)
    add_textbox(slide, Inches(3.3), y + Inches(0.3), Inches(3.3), Inches(0.25),
                reasoning, 10, GRAY)

# Right: Revenue
add_textbox(slide, Inches(7.2), Inches(1.3), Inches(5.5), Inches(0.5),
            "Revenue Model", 20, TEAL, bold=True)

rev_cards = [
    ("Freemium", "Basic insights: FREE\nPremium (AI coaching, history): ₹99/mo", TEAL),
    ("B2B Wellness", "Employer partnerships\n₹50/employee/month", PURPLE),
    ("Affiliate", "Commission-free FD, MF,\ninsurance product referrals", GOLD),
    ("API License", "Behavioral engine as a service\nfor neobanks & NBFCs", ERROR),
]

for i, (title, desc, accent) in enumerate(rev_cards):
    y = Inches(1.9 + i * 1.2)
    add_card(slide, Inches(7.2), y, Inches(5.3), Inches(1.0))
    add_textbox(slide, Inches(7.5), y + Inches(0.1), Inches(1.5), Inches(0.4),
                title, 16, accent, bold=True)
    add_textbox(slide, Inches(9.2), y + Inches(0.1), Inches(3.1), Inches(0.8),
                desc, 12, WHITE)

# ════════════════════════════════════════════════════════════════
# SLIDE 11 — Competitive Moat
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "What Makes FintLer Defensible", 32, WHITE, bold=True)

moats = [
    ("Data Moat", "Gmail integration captures ALL bank transactions — not just partner banks. The user's full financial picture in one place.", TEAL),
    ("AI Moat", "Gemini 2.0 Flash parsing trained on Indian bank emails. 95%+ accuracy. Regex-first → AI-fallback pipeline. Gets smarter with every user feedback loop.", PURPLE),
    ("Behavioral Moat", "Personality profiling + micro-goals create habit loops. No competitor offers 'Spending Personality + Weekly Goal' as a compound feature.", GOLD),
    ("Trust Moat", "Read-only consent, zero-judgment tone, privacy-first design. Users feel safe — not shamed. This drives retention and word-of-mouth.", ERROR),
    ("Distribution Moat", "Gmail = universal. Every Indian salary earner has one. No app download required for initial engagement. OAuth is the lowest-friction onboarding.", TEAL),
]

for i, (title, desc, accent) in enumerate(moats):
    y = Inches(1.4 + i * 1.1)
    add_card(slide, Inches(0.8), y, Inches(2.5), Inches(0.9))
    add_textbox(slide, Inches(1.0), y + Inches(0.15), Inches(2.1), Inches(0.6),
                title, 18, accent, bold=True)
    add_card(slide, Inches(3.5), y, Inches(8.8), Inches(0.9))
    add_textbox(slide, Inches(3.7), y + Inches(0.15), Inches(8.4), Inches(0.6),
                desc, 13, GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 12 — Team & Current Status
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
            "Current Status & Achievements", 32, WHITE, bold=True)

achievements = [
    ("✅ Phase 1: Security", "Hardened CORS, Auth, Schema — production-ready security posture"),
    ("✅ Phase 2: Error Handling", "ErrorBoundary, modal controls, logger — crash-proof UX"),
    ("✅ Phase 3: Schema Migration", "Consolidated Supabase schema with RLS, realtime, email_connections"),
    ("✅ Phase 4: Frontend Cleanup", "Code splitting, mobile navbar, memo components — sub-1s loads"),
    ("✅ Phase 5: Edge Functions", "6 functions optimized: JWT auth, Gemini 2.0, consistent error handling"),
    ("✅ Phase 6: Testing", "19 unit tests (Vitest) — logger, components, routing, error boundary"),
    ("✅ Phase 7: Production", "Deployed on Netlify — fintlerai.netlify.app — Supabase linked & configured"),
]

for i, (title, desc) in enumerate(achievements):
    y = Inches(1.4 + i * 0.75)
    add_textbox(slide, Inches(0.8), y, Inches(3.5), Inches(0.5),
                title, 14, TEAL, bold=True)
    add_textbox(slide, Inches(4.3), y, Inches(8.0), Inches(0.5),
                desc, 13, GRAY)

add_card(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8))
add_textbox(slide, Inches(1.0), Inches(6.3), Inches(11.0), Inches(0.6),
            "Live Preview: https://fintlerai.netlify.app  |  Supabase: dpaygeohgnrmwucsbcak.supabase.co  |  Netlify CI: GitHub auto-deploy on main", 14, TEAL)

# ════════════════════════════════════════════════════════════════
# SLIDE 13 — Ask / Next Steps
# ════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_orb(slide, Inches(0), Inches(1), Inches(7), TEAL, 0.08)
add_orb(slide, Inches(8), Inches(-1), Inches(6), PURPLE, 0.08)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.8),
            "Thank You", 48, WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(0.6),
            "Let's cure salary day anxiety — together.", 24, TEAL)
add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.5),
            "Live Demo: https://fintlerai.netlify.app", 16, GRAY)
add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11), Inches(0.5),
            "Product Analysis: C:\\Product Sprint 2\\PRODUCT_ANALYSIS.md", 14, RGBColor(0x60, 0x60, 0x70))

# ── Save ──
output_path = "C:\\Product Sprint 2\\FintLer_Product_Deck.pptx"
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
